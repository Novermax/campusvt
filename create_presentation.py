"""
Script per generare la presentazione tecnica Campus Virtual Training
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === COLORI TEMA ===
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_MEDIUM = RGBColor(0x16, 0x21, 0x3e)
BG_CARD = RGBColor(0x1f, 0x2b, 0x4d)
ACCENT_BLUE = RGBColor(0x00, 0xd2, 0xff)
ACCENT_GREEN = RGBColor(0x00, 0xe6, 0x96)
ACCENT_ORANGE = RGBColor(0xff, 0x9f, 0x43)
ACCENT_PURPLE = RGBColor(0xa2, 0x9b, 0xfe)
ACCENT_RED = RGBColor(0xff, 0x6b, 0x6b)
ACCENT_YELLOW = RGBColor(0xff, 0xd9, 0x3d)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xb0, 0xb0, 0xc0)
MID_GRAY = RGBColor(0x80, 0x80, 0x99)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, default_size=16, default_color=LIGHT_GRAY, font_name='Segoe UI'):
    """lines = list of (text, size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            line_data = (line_data, default_size, default_color, False, PP_ALIGN.LEFT)
        text, size, color, bold, align = line_data
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(4)
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    # Adjust rounding
    shape.adjustments[0] = 0.05
    return shape

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1 - TITOLO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Linea accent top
add_accent_line(slide, 0, 0, 13.333, ACCENT_BLUE)

# Titolo
add_text_box(slide, 1.5, 1.5, 10, 1.2, 'CAMPUS VIRTUAL TRAINING', 44, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 2.6, 10, 0.8, 'Panoramica Tecnica per l\'Ufficio Tecnico', 24, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Linea separatrice
add_accent_line(slide, 4.5, 3.6, 4.3, ACCENT_BLUE)

# Bullet points
bullets = [
    'Simulatore 3D di procedure di manutenzione industriale',
    'Browser-based  |  Zero installazioni  |  Desktop + Mobile',
    'Procedure interamente data-driven (file .cvtscript)',
]
for i, b in enumerate(bullets):
    add_text_box(slide, 2.5, 4.3 + i * 0.55, 8, 0.5, b, 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Footer
add_text_box(slide, 1, 6.5, 11, 0.5, 'Three.js r155  |  Vanilla ES6  |  CSS Modulare  |  ~370.000 righe', 14, MID_GRAY, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 - STACK TECNOLOGICO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, 0, 0, 13.333, ACCENT_GREEN)

add_text_box(slide, 0.8, 0.3, 11, 0.8, 'Stack Tecnologico', 32, WHITE, True)
add_accent_line(slide, 0.8, 1.0, 2.5, ACCENT_GREEN)

# Cards
cards = [
    ('Rendering 3D', 'Three.js r155 (ES Modules)', 'WebGL, raycasting, GLTF/DRACO/OBJ/STL\nLoaders nativi Three.js\nNessun wrapper custom', ACCENT_BLUE),
    ('Frontend', 'Vanilla ES6 — No Framework', 'Nessun React/Vue/Angular\nNessun bundler (Webpack/Vite)\nImport nativi del browser', ACCENT_GREEN),
    ('CSS', '4 File Modulari', 'base.css — reset e fondamenta\ncomponents.css — UI riutilizzabili\nlayout.css — posizionamento\npages.css — stili per pagina', ACCENT_ORANGE),
    ('Server', 'File Server Statico', 'python -m http.server 8000\nNessun backend, nessuna API\nNessun database', ACCENT_PURPLE),
    ('Target', 'Cross-Browser + Mobile', 'Chrome 80+, Firefox 75+\nSafari 13+, Edge 80+\nTouch nativo con gesture dedicate', ACCENT_RED),
]

for i, (title, subtitle, desc, color) in enumerate(cards):
    col = i % 3
    row = i // 3
    x = 0.6 + col * 4.1
    y = 1.5 + row * 2.9
    add_rounded_rect(slide, x, y, 3.8, 2.5, BG_CARD, color)
    add_text_box(slide, x + 0.25, y + 0.2, 3.3, 0.4, title, 18, color, True)
    add_text_box(slide, x + 0.25, y + 0.6, 3.3, 0.4, subtitle, 13, WHITE, True)
    add_multiline_box(slide, x + 0.25, y + 1.0, 3.3, 1.4,
                      [(line, 12, LIGHT_GRAY, False, PP_ALIGN.LEFT) for line in desc.split('\n')])

# Nota bottom
add_text_box(slide, 0.8, 6.8, 12, 0.4, 'Zero dipendenze build-time: si edita il file, si refresha il browser.', 14, ACCENT_YELLOW, False)


# ============================================================
# SLIDE 3 - ARCHITETTURA GENERALE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, 0, 0, 13.333, ACCENT_BLUE)

add_text_box(slide, 0.8, 0.3, 11, 0.8, 'Architettura Generale', 32, WHITE, True)
add_accent_line(slide, 0.8, 1.0, 2.8, ACCENT_BLUE)

# Bootstrap box
add_rounded_rect(slide, 0.6, 1.5, 12.1, 1.0, BG_CARD, ACCENT_BLUE)
add_text_box(slide, 0.8, 1.55, 2, 0.4, 'BOOTSTRAP', 14, ACCENT_BLUE, True)
add_text_box(slide, 0.8, 1.9, 11.5, 0.5, 'index.html  \u2192  app.js  \u2192  THREE.js (import ES Module)  \u2192  Scene3D.init()  \u2192  UI.init()  \u2192  Sistemi Core', 13, LIGHT_GRAY, False)

# Scene3D column
add_rounded_rect(slide, 0.6, 2.8, 3.7, 3.8, BG_CARD, ACCENT_BLUE)
add_text_box(slide, 0.8, 2.9, 3.3, 0.4, 'Scene3D', 20, ACCENT_BLUE, True)
add_text_box(slide, 0.8, 3.25, 3.3, 0.35, 'scene3d-modular.js  (5.250 righe)', 11, MID_GRAY, False)
modules_3d = [
    'CameraControls.js — orbit, zoom, transizioni',
    'AnimationSystem.js — traslazioni, rotazioni',
    'MultiStepAnimationSystem.js — catene N azioni',
    'HighlightSystem.js — feedback visivo',
    'MovementParser.js — parsing comandi azione',
    'ParticleSystem.js — effetti particellari',
]
for j, m in enumerate(modules_3d):
    add_text_box(slide, 1.0, 3.7 + j * 0.42, 3.2, 0.4, m, 10, LIGHT_GRAY, False)

# UI column
add_rounded_rect(slide, 4.5, 2.8, 3.7, 3.8, BG_CARD, ACCENT_GREEN)
add_text_box(slide, 4.7, 2.9, 3.3, 0.4, 'UI', 20, ACCENT_GREEN, True)
add_text_box(slide, 4.7, 3.25, 3.3, 0.35, 'ui.js  (6.100 righe)', 11, MID_GRAY, False)
modules_ui = [
    'TutorialManager.js — parsing .cvtscript',
    'ScenarioManager.js — caricamento scenari',
    'ToolsManager.js — toolbar + cursori SVG',
    'PageManager.js — routing Home/Scenario',
    'DynamicToolStyles.js — CSS runtime',
    'ModelManager.js — gestione modelli',
]
for j, m in enumerate(modules_ui):
    add_text_box(slide, 4.9, 3.7 + j * 0.42, 3.2, 0.4, m, 10, LIGHT_GRAY, False)

# Core Systems column
add_rounded_rect(slide, 8.4, 2.8, 4.3, 3.8, BG_CARD, ACCENT_ORANGE)
add_text_box(slide, 8.6, 2.9, 3.9, 0.4, 'Sistemi Core', 20, ACCENT_ORANGE, True)
add_text_box(slide, 8.6, 3.25, 3.9, 0.35, 'js/core/  (23 moduli indipendenti)', 11, MID_GRAY, False)
modules_core = [
    'DragDropSystem + SnapSystem — drag & snap',
    'ScreenSystem — schermi 2D nei modelli 3D',
    'HoldableSystem — oggetti impugnabili',
    'InteractiveObject3D — pulsanti, LED, rotary',
    'StepController — trigger multi-sorgente',
    'StepGatingManager — blocco/sblocco pulsanti',
    'AnimatedWindowSystem — sequenze immagini',
    'ToolRegistry — registry utensili per scenario',
]
for j, m in enumerate(modules_core):
    add_text_box(slide, 8.8, 3.7 + j * 0.38, 3.8, 0.38, m, 10, LIGHT_GRAY, False)

# Comunicazione
add_text_box(slide, 0.6, 6.8, 12, 0.4, 'Comunicazione tra moduli: window.* globals (Scene3D, UI, DragDropSystem, ...) — semplice ed efficace per app monolitica', 12, MID_GRAY, False)


# ============================================================
# SLIDE 4 - FORMATO .cvtscript
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, 0, 0, 13.333, ACCENT_ORANGE)

add_text_box(slide, 0.8, 0.3, 11, 0.8, 'Il Formato .cvtscript — Data-Driven', 32, WHITE, True)
add_accent_line(slide, 0.8, 1.0, 3.6, ACCENT_ORANGE)

add_text_box(slide, 0.8, 1.3, 11, 0.5, 'Tutta la logica delle procedure e\u0300 definita in file di configurazione .cvtscript (sintassi INI-like). Zero codice JS per nuovi scenari.', 15, LIGHT_GRAY, False)

# 3 livelli
levels = [
    ('1', 'home_config.cvtscript', 'Indice globale degli scenari\n+ direzioni svitatura per modello', ACCENT_BLUE,
     '[Manutenzione_Elettromandrino]\nScenario=scenes/Manutenzione_Elettromandrino\nTutorial=scenes/.../tutorial.cvtscript\n\n[models/vite.glb]\ndirection=0,0,1'),
    ('2', 'config.cvtscript (per scenario)', 'Utensili disponibili,\nicone, cursori SVG custom', ACCENT_GREEN,
     '[Tools]\nTool=brugola\nTool=chiave_inglese\n\n[Tool:brugola]\nLabel=Chiave a Brugola\nIcon=utilimages/brugola.png\nCursor=cursors/brugola_normale.svg'),
    ('3', 'tutorial.cvtscript (per scenario)', 'Procedura step-by-step\ncompleta', ACCENT_ORANGE,
     '[Step 3 - Rimuovi vite coperchio]\nElemento=models/vite_coperchio.glb\nUtensile=ChiaveBrugola\nAzione1=svita(0.5)\nDescrizione=Svitare la vite\nCameraDistance=0.8'),
]

for i, (num, title, desc, color, code) in enumerate(levels):
    x = 0.6 + i * 4.2
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.0), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, x + 0.6, 2.0, 3.2, 0.5, title, 14, color, True)
    add_text_box(slide, x + 0.6, 2.4, 3.2, 0.6, desc, 11, LIGHT_GRAY, False)

    # Code box
    add_rounded_rect(slide, x, 3.2, 3.9, 3.2, RGBColor(0x0d, 0x11, 0x20), MID_GRAY)
    add_multiline_box(slide, x + 0.15, 3.35, 3.6, 3.0,
                      [(line, 10, ACCENT_GREEN if line.startswith('[') else (MID_GRAY if line.startswith('#') else LIGHT_GRAY),
                        line.startswith('['), PP_ALIGN.LEFT) for line in code.split('\n')],
                      font_name='Consolas')

add_text_box(slide, 0.8, 6.7, 12, 0.4, 'Nuovo scenario = nuova cartella con modelli .glb + 2 file .cvtscript. Nessuna riga di JavaScript.', 14, ACCENT_YELLOW, True)


# ============================================================
# SLIDE 5 - FLUSSO RUNTIME DI UNO STEP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, 0, 0, 13.333, ACCENT_PURPLE)

add_text_box(slide, 0.8, 0.3, 11, 0.8, 'Flusso Runtime di uno Step', 32, WHITE, True)
add_accent_line(slide, 0.8, 1.0, 3.2, ACCENT_PURPLE)

steps = [
    ('1', 'UI.goToStep(N)', 'Legge la config dello step\ndal tutorial parsato', ACCENT_BLUE),
    ('2', 'Camera Transition', 'Transizione animata verso\nCameraPos/Target/Distance', ACCENT_GREEN),
    ('3', 'Highlight Target', 'Emissive giallo pulsante\nsull\'elemento target', ACCENT_YELLOW),
    ('4', 'Mostra Fumetto', 'Descrizione step +\nicona utensile richiesto', ACCENT_ORANGE),
    ('5', 'Selezione Utensile', 'Utente sceglie dalla toolbar\nCursore SVG custom', ACCENT_PURPLE),
    ('6', 'Click Elemento', 'handleModelAction() verifica\nutensile + target corretti', ACCENT_RED),
    ('7', 'Animazione', 'AnimationSystem esegue\nAzione1, Azione2, ... in sequenza', ACCENT_BLUE),
    ('8', 'DrivenObjects', 'Oggetti slave si muovono\nin parallelo al master', ACCENT_GREEN),
    ('9', 'Avanzamento', 'Completamento animazione\n\u2192 step successivo', ACCENT_YELLOW),
]

for i, (num, title, desc, color) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.2
    y = 1.5 + row * 2.0

    # Step card
    add_rounded_rect(slide, x, y, 3.8, 1.6, BG_CARD, color)

    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.15), Inches(y + 0.15), Inches(0.4), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = BG_DARK
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + 0.65, y + 0.1, 3.0, 0.35, title, 14, WHITE, True)
    add_text_box(slide, x + 0.15, y + 0.6, 3.5, 0.9, desc, 11, LIGHT_GRAY, False)

    # Arrow between cards (same row)
    if col < 2:
        arrow_x = x + 3.85
        arrow_y = y + 0.7
        add_text_box(slide, arrow_x, arrow_y, 0.4, 0.3, '\u2192', 18, MID_GRAY, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6 - SISTEMA TOUCH
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, 0, 0, 13.333, ACCENT_RED)

add_text_box(slide, 0.8, 0.3, 11, 0.8, 'Sistema Touch Mobile', 32, WHITE, True)
add_accent_line(slide, 0.8, 1.0, 2.8, ACCENT_RED)

# Left: Architecture
add_rounded_rect(slide, 0.6, 1.5, 6.0, 2.5, BG_CARD, ACCENT_RED)
add_text_box(slide, 0.8, 1.6, 5.5, 0.4, 'Architettura: 8 moduli in js/touch/', 16, ACCENT_RED, True)
touch_modules = [
    'index.js — coordinatore principale',
    'TouchEventDispatcher — normalizza eventi touch',
    'GestureRecognizer — riconosce tap, drag, pinch',
    'TouchInputRouter — routing a layer con priorita\u0300',
    'TouchCameraHandler — orbit + zoom (pinch 2 dita)',
    'TouchDragHandler — drag&drop + ghost mesh',
    'TouchUIHandler — interazione UI overlay',
    'TouchInteractive3DHandler — pulsanti 3D',
]
for j, m in enumerate(touch_modules):
    add_text_box(slide, 1.0, 2.1 + j * 0.32, 5.3, 0.3, m, 10, LIGHT_GRAY, False)

# Right: Priority
add_rounded_rect(slide, 6.9, 1.5, 5.8, 2.5, BG_CARD, ACCENT_ORANGE)
add_text_box(slide, 7.1, 1.6, 5.3, 0.4, 'Priorita\u0300 Routing (determineLayer)', 16, ACCENT_ORANGE, True)

priorities = [
    ('UI', '3', 'Overlay, bottoni HTML', ACCENT_RED),
    ('Interactive3D', '2', 'Pulsanti 3D, LED, rotary', ACCENT_ORANGE),
    ('Object3D', '1', 'Modelli trascinabili', ACCENT_YELLOW),
    ('Camera', '0', 'Orbit, zoom (fallback)', ACCENT_GREEN),
]
for j, (name, pri, desc, color) in enumerate(priorities):
    y = 2.2 + j * 0.48
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(y), Inches(0.35), Inches(0.3))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = pri
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = BG_DARK
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text_box(slide, 7.65, y, 1.8, 0.3, name, 12, color, True)
    add_text_box(slide, 9.5, y, 3.0, 0.3, desc, 11, LIGHT_GRAY, False)

# Bottom: Key rule
add_rounded_rect(slide, 0.6, 4.3, 12.1, 1.4, BG_CARD, ACCENT_YELLOW)
add_text_box(slide, 0.8, 4.4, 11.5, 0.4, 'REGOLA FONDAMENTALE: Un solo evento per tap', 18, ACCENT_YELLOW, True)
add_multiline_box(slide, 0.8, 4.9, 11.5, 0.9, [
    ('Tap su elemento con azione tutorial  \u2192  SOLO esecuzione azione (NO pivot camera)', 13, ACCENT_GREEN, False, PP_ALIGN.LEFT),
    ('Tap su oggetto generico  \u2192  SOLO pivot camera (NO azione)', 13, ACCENT_BLUE, False, PP_ALIGN.LEFT),
    ('Desktop e Touch sono sistemi separati: handler mouse si disabilitano quando TouchSystem e\u0300 attivo', 12, MID_GRAY, False, PP_ALIGN.LEFT),
])

# Gesture summary
add_rounded_rect(slide, 0.6, 5.9, 12.1, 1.3, BG_CARD, ACCENT_PURPLE)
add_text_box(slide, 0.8, 6.0, 5, 0.35, 'Gesture Attive', 16, ACCENT_PURPLE, True)
gestures = [
    ('1 dito tap', 'Azione tool OPPURE pivot camera (mai entrambi)'),
    ('1 dito drag', 'Camera orbit (layer=CAMERA) oppure DragDrop (se step ha DragDrop=true)'),
    ('2 dita pinch', 'Zoom (sensitivity: 0.003)'),
    ('2 dita drag / double-tap', 'DISABILITATI'),
]
for j, (gesture, desc) in enumerate(gestures):
    add_text_box(slide, 1.0, 6.35 + j * 0.25, 2.0, 0.25, gesture, 10, ACCENT_PURPLE, True)
    add_text_box(slide, 3.2, 6.35 + j * 0.25, 9.0, 0.25, desc, 10, LIGHT_GRAY, False)


# ============================================================
# SLIDE 7 - DRAG & DROP + SNAP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, 0, 0, 13.333, ACCENT_GREEN)

add_text_box(slide, 0.8, 0.3, 11, 0.8, 'DragDrop + Snap System', 32, WHITE, True)
add_accent_line(slide, 0.8, 1.0, 3.0, ACCENT_GREEN)

# Snap modes
snap_modes = [
    ('Coordinate Fisse', 'Globali o per-oggetto\nSnapPoint=(0.5,0.2,0.3)', ACCENT_BLUE,
     'Tutti gli oggetti snappano\nai punti definiti.\nAuto-detect formato:\nassenza ":" = globale'),
    ('Pivot Point', 'Snap al pivot, non al centro BB\nSnapPoint=pivot:(x,y,z)', ACCENT_GREEN,
     'Usa il pivot point del\nmodello invece del centro\ndel bounding box'),
    ('Offset', 'Dalla posizione originale\nSnapPoint=offset:(0,0,0.02)', ACCENT_ORANGE,
     'Snap relativo alla\nposizione di caricamento\n+ offset specificato'),
    ('SnapTargets', 'Posizione originale di altri modelli\nSnapTargets=foro_1_original', ACCENT_PURPLE,
     'Snap alla posizione dove\nun altro modello e\u0300 stato\ncaricato (es. vite \u2192 foro)'),
]

for i, (title, syntax, color, desc) in enumerate(snap_modes):
    x = 0.6 + i * 3.15
    add_rounded_rect(slide, x, 1.5, 2.95, 3.5, BG_CARD, color)
    add_text_box(slide, x + 0.15, 1.6, 2.65, 0.35, title, 15, color, True)

    # Syntax box
    add_rounded_rect(slide, x + 0.1, 2.05, 2.75, 0.7, RGBColor(0x0d, 0x11, 0x20))
    add_multiline_box(slide, x + 0.2, 2.1, 2.55, 0.6,
                      [(line, 9, ACCENT_GREEN, False, PP_ALIGN.LEFT) for line in syntax.split('\n')],
                      font_name='Consolas')

    add_multiline_box(slide, x + 0.15, 2.85, 2.65, 1.8,
                      [(line, 10, LIGHT_GRAY, False, PP_ALIGN.LEFT) for line in desc.split('\n')])

# Ghost mesh section
add_rounded_rect(slide, 0.6, 5.3, 12.1, 1.0, BG_CARD, ACCENT_YELLOW)
add_text_box(slide, 0.8, 5.4, 3, 0.35, 'Ghost Mesh (Touch)', 16, ACCENT_YELLOW, True)
add_text_box(slide, 0.8, 5.75, 11.5, 0.5, 'Mesh semi-trasparente segue il dito in real-time durante il drag.  Verde = snap raggiungibile  |  Rosso = fuori range.  Soglia: DragDropDistance (default 0.3)', 12, LIGHT_GRAY, False)

# Reference _original
add_rounded_rect(slide, 0.6, 6.5, 12.1, 0.7, BG_CARD, MID_GRAY)
add_text_box(slide, 0.8, 6.55, 11.5, 0.35, 'Riferimento _original', 14, WHITE, True)
add_text_box(slide, 0.8, 6.85, 11.5, 0.3, 'modello_original e\u0300 un riferimento virtuale creato automaticamente alla posizione di caricamento. Usato in SnapTargets e traslazioni relative.', 11, LIGHT_GRAY, False)


# ============================================================
# SLIDE 8 - SISTEMI INTERATTIVI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, 0, 0, 13.333, ACCENT_PURPLE)

add_text_box(slide, 0.8, 0.3, 11, 0.8, 'Sistemi Interattivi Avanzati', 32, WHITE, True)
add_accent_line(slide, 0.8, 1.0, 3.2, ACCENT_PURPLE)

# ScreenSystem
add_rounded_rect(slide, 0.6, 1.5, 4.0, 3.0, BG_CARD, ACCENT_BLUE)
add_text_box(slide, 0.8, 1.6, 3.5, 0.35, 'ScreenSystem', 18, ACCENT_BLUE, True)
add_text_box(slide, 0.8, 1.95, 3.5, 0.3, 'Schermi 2D nei modelli 3D', 11, MID_GRAY, False)
screen_items = [
    'Schermi con viste multiple (GLB separati)',
    'Hotspot cliccabili (posizione, size, label)',
    'Navigazione tra viste (NextView=)',
    'Azioni su click (animazioni, suoni)',
    'Camera perpendicolare automatica',
    'Simula pannelli controllo CNC',
]
for j, item in enumerate(screen_items):
    add_text_box(slide, 0.9, 2.35 + j * 0.33, 3.5, 0.3, '\u2022 ' + item, 10, LIGHT_GRAY, False)

# InteractiveObject3D
add_rounded_rect(slide, 4.8, 1.5, 4.0, 3.0, BG_CARD, ACCENT_ORANGE)
add_text_box(slide, 5.0, 1.6, 3.5, 0.35, 'InteractiveObject3D', 18, ACCENT_ORANGE, True)
add_text_box(slide, 5.0, 1.95, 3.5, 0.3, 'Oggetti complessi con figli interattivi', 11, MID_GRAY, False)
io_items = [
    ('button', 'click \u2192 evento'),
    ('rotary', 'click \u2192 cicla stati + rotazione'),
    ('indicator', 'visibilita\u0300 da stato (LED on/off)'),
    ('screen', 'visibilita\u0300 da currentScreen'),
]
for j, (tipo, desc) in enumerate(io_items):
    add_text_box(slide, 5.1, 2.4 + j * 0.38, 1.2, 0.35, tipo, 11, ACCENT_ORANGE, True, font_name='Consolas')
    add_text_box(slide, 6.3, 2.4 + j * 0.38, 2.3, 0.35, desc, 11, LIGHT_GRAY, False)
add_text_box(slide, 5.1, 4.0, 3.5, 0.3, 'StateGroup: varianti mutuamente esclusive', 10, MID_GRAY, False)

# StepController + Gating
add_rounded_rect(slide, 9.0, 1.5, 3.7, 3.0, BG_CARD, ACCENT_GREEN)
add_text_box(slide, 9.2, 1.6, 3.3, 0.35, 'StepController', 18, ACCENT_GREEN, True)
add_text_box(slide, 9.2, 1.95, 3.3, 0.3, 'Trigger da sorgenti diverse', 11, MID_GRAY, False)
trigger_items = [
    ('Physical', 'Pulsante 3D (mesh click)'),
    ('Screen', 'Hotspot su schermo'),
    ('Holdable', 'Oggetto in mano (use)'),
]
for j, (src, desc) in enumerate(trigger_items):
    add_text_box(slide, 9.3, 2.4 + j * 0.35, 1.5, 0.3, src, 11, ACCENT_GREEN, True)
    add_text_box(slide, 10.6, 2.4 + j * 0.35, 2.0, 0.3, desc, 11, LIGHT_GRAY, False)

add_text_box(slide, 9.3, 3.5, 3.3, 0.5, 'StepGatingManager:\nstepIndex = -1 \u2192 tutti i pulsanti bloccati. ActiveButtons= abilita selettivamente.', 10, MID_GRAY, False)

# HoldableSystem + AnimatedWindow
add_rounded_rect(slide, 0.6, 4.8, 6.1, 1.6, BG_CARD, ACCENT_PURPLE)
add_text_box(slide, 0.8, 4.9, 5.5, 0.35, 'HoldableSystem — Oggetti Impugnabili', 16, ACCENT_PURPLE, True)
add_multiline_box(slide, 0.8, 5.3, 5.5, 1.0, [
    ('HoldAction=pick  \u2192  prende oggetto (attaccato alla camera)', 11, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ('HoldAction=release  \u2192  rilascia oggetto', 11, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ('HoldState=held  \u2192  step richiede oggetto gia\u0300 in mano', 11, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ('Usato per: telecomandi, strumenti di misura, ecc.', 11, MID_GRAY, False, PP_ALIGN.LEFT),
])

add_rounded_rect(slide, 6.9, 4.8, 5.8, 1.6, BG_CARD, ACCENT_YELLOW)
add_text_box(slide, 7.1, 4.9, 5.3, 0.35, 'AnimatedWindowSystem — Sequenze Immagini', 16, ACCENT_YELLOW, True)
add_multiline_box(slide, 7.1, 5.3, 5.3, 1.0, [
    ('Overlay 2D con sequenza immagini (frame-by-frame)', 11, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ('Posizionamento: center, top-left, coordinate custom', 11, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ('MaxTriggers = N trigger avanti+indietro prima di chiusura', 11, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ('Usato per: istruzioni visive, diagrammi, schemi', 11, MID_GRAY, False, PP_ALIGN.LEFT),
])

# AutoExecute
add_rounded_rect(slide, 0.6, 6.6, 12.1, 0.7, BG_CARD, ACCENT_RED)
add_text_box(slide, 0.8, 6.65, 11.5, 0.3, 'AutoExecute / Autoaction — Step Automatici', 14, ACCENT_RED, True)
add_text_box(slide, 0.8, 6.95, 11.5, 0.3, 'AutoExecute=true \u2192 animazione automatica a T+300ms  |  Autoaction=true \u2192 auto-equip utensile + auto-esegui + auto-avanza  |  Con DragDrop: auto-snap senza drag utente', 11, LIGHT_GRAY, False)


# ============================================================
# SLIDE 9 - AZIONI ANIMATE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, 0, 0, 13.333, ACCENT_YELLOW)

add_text_box(slide, 0.8, 0.3, 11, 0.8, 'Azioni Animate — Riferimento Comandi', 32, WHITE, True)
add_accent_line(slide, 0.8, 1.0, 3.8, ACCENT_YELLOW)

actions = [
    ('traslazione:(x,y,z,durata)', 'Movimento lineare a coordinate', ACCENT_BLUE),
    ('rotazione:(rx,ry,rz,durata)', 'Rotazione (in gradi)', ACCENT_BLUE),
    ('appoggia(durata)', 'Posiziona a Y=0 (calcola BB automaticamente)', ACCENT_BLUE),
    ('svita / svita(0.8)', 'Rotazione + traslazione lungo direction (da home_config)', ACCENT_GREEN),
    ('avvita / avvita(0.8)', 'Inverso di svita', ACCENT_GREEN),
    ('estrai / estrai(0.6)', 'Traslazione lungo direction', ACCENT_ORANGE),
    ('inserisci / inserisci(0.6)', 'Inverso di estrai', ACCENT_ORANGE),
    ('oscillazione:ref,(pA,pB,cicli,durata)', 'Movimento alternato ripetitivo (es. pompaggio)', ACCENT_PURPLE),
]

# Header
add_rounded_rect(slide, 0.6, 1.4, 12.1, 0.5, BG_CARD)
add_text_box(slide, 0.8, 1.42, 5, 0.4, 'Comando', 13, MID_GRAY, True)
add_text_box(slide, 6.0, 1.42, 6, 0.4, 'Descrizione', 13, MID_GRAY, True)

for i, (cmd, desc, color) in enumerate(actions):
    y = 2.0 + i * 0.48
    if i % 2 == 0:
        add_rounded_rect(slide, 0.6, y, 12.1, 0.45, RGBColor(0x14, 0x1c, 0x38))
    add_text_box(slide, 0.8, y + 0.02, 5, 0.4, cmd, 12, color, False, font_name='Consolas')
    add_text_box(slide, 6.0, y + 0.02, 6, 0.4, desc, 12, LIGHT_GRAY, False)

# DrivenObjects
add_rounded_rect(slide, 0.6, 5.9, 12.1, 1.3, BG_CARD, ACCENT_GREEN)
add_text_box(slide, 0.8, 6.0, 5, 0.35, 'DrivenObjects — Movimenti Paralleli', 16, ACCENT_GREEN, True)
add_multiline_box(slide, 0.8, 6.4, 11.5, 0.8, [
    ('DrivenObjects=flangia.glb,traslazione:(0,0,0.1,0.5)     \u2192  azione indipendente in parallelo al master', 11, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ('DrivenObjects=tubograsso.glb,follow                       \u2192  segue rigidamente il master 1:1', 11, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ('Solo il master controlla l\'avanzamento. I driven non bloccano il tutorial.', 11, MID_GRAY, False, PP_ALIGN.LEFT),
])


# ============================================================
# SLIDE 10 - DEBUG E FAST-FORWARD
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, 0, 0, 13.333, ACCENT_BLUE)

add_text_box(slide, 0.8, 0.3, 11, 0.8, 'Debug Console + Fast-Forward', 32, WHITE, True)
add_accent_line(slide, 0.8, 1.0, 3.4, ACCENT_BLUE)

# Debug commands
debug_sections = [
    ('Camera', [
        'Scene3D.getCameraInfo()  — posizione + sintassi pronta per tutorial',
        'Scene3D.setCameraFromInfo({...})  — imposta camera programmaticamente',
        'Scene3D.exportCurrentModelPositions()  — download Posizione=/Rotazione=',
        'Scene3D.listChildNames(model)  — nomi per TargetChild',
    ], ACCENT_BLUE),
    ('Navigazione Tutorial', [
        'jumpToStep(5)  — salta con fast-forward (applica trasformazioni precedenti)',
        'jumpToStep(5, false)  — salta senza fast-forward (solo debug)',
        'listSteps()  — lista tutti gli step',
        'findStep("vite")  — cerca per keyword, auto-jump se 1 risultato',
    ], ACCENT_GREEN),
    ('Sistemi Interattivi', [
        'DragDropSystem.debugSnapSystem()  — stato snap points',
        'ScreenSystem.debugInfo()  — stato schermi',
        'InteractiveObject3D.getState("pulpito")  — stato corrente',
        'StepController.simulateTrigger("physical", "pulpito.Btn")  — simula input',
    ], ACCENT_ORANGE),
]

for i, (title, cmds, color) in enumerate(debug_sections):
    x = 0.6 + i * 4.2
    add_rounded_rect(slide, x, 1.4, 3.9, 3.4, BG_CARD, color)
    add_text_box(slide, x + 0.15, 1.5, 3.6, 0.35, title, 15, color, True)
    for j, cmd in enumerate(cmds):
        parts = cmd.split('  — ')
        add_text_box(slide, x + 0.15, 1.95 + j * 0.7, 3.6, 0.4, parts[0], 9, WHITE, False, font_name='Consolas')
        if len(parts) > 1:
            add_text_box(slide, x + 0.15, 2.25 + j * 0.7, 3.6, 0.3, parts[1], 9, MID_GRAY, False)

# Fast-forward explanation
add_rounded_rect(slide, 0.6, 5.1, 12.1, 2.1, BG_CARD, ACCENT_YELLOW)
add_text_box(slide, 0.8, 5.2, 5, 0.4, 'Fast-Forward: jumpToStep(N)', 18, ACCENT_YELLOW, True)
add_multiline_box(slide, 0.8, 5.6, 11.5, 1.5, [
    ('Quando si salta a step N, il sistema applica retroattivamente tutte le trasformazioni degli step 0..N-1', 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ('in modo istantaneo. La scena risulta coerente come se l\'utente avesse completato ogni step.', 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ('', 8, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ('Trasformazioni applicate: Posizione, Rotazione, svita, avvita, estrai, inserisci, traslazione', 12, ACCENT_GREEN, False, PP_ALIGN.LEFT),
    ('Timeout sicurezza: 5 secondi max per step, poi avanza comunque', 12, ACCENT_ORANGE, False, PP_ALIGN.LEFT),
    ('Ereditarieta\u0300 camera: le proprieta\u0300 camera si ereditano dal contesto precedente (sezione o step)', 12, ACCENT_BLUE, False, PP_ALIGN.LEFT),
])


# ============================================================
# SLIDE 11 - RIEPILOGO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_accent_line(slide, 0, 0, 13.333, ACCENT_BLUE)

add_text_box(slide, 1.5, 0.6, 10, 0.8, 'Riepilogo', 36, WHITE, True, PP_ALIGN.CENTER)
add_accent_line(slide, 5.5, 1.3, 2.3, ACCENT_BLUE)

summary_items = [
    ('Browser-native', 'Three.js + Vanilla ES6, nessun framework, nessun build step', ACCENT_BLUE),
    ('Data-driven', 'Nuovi scenari creati interamente con file .cvtscript, zero codice', ACCENT_GREEN),
    ('Modulare', '23+ moduli core indipendenti, ciascuno con API globale esposta', ACCENT_ORANGE),
    ('Cross-platform', 'Desktop (mouse) + Mobile (touch) con sistemi completamente separati', ACCENT_PURPLE),
    ('Interattivo', 'Schermi, pulsanti 3D, LED, drag&drop con snap magnetico', ACCENT_RED),
    ('Debug-friendly', 'Console ricca per iterazione rapida su posizioni camera e step', ACCENT_YELLOW),
]

for i, (title, desc, color) in enumerate(summary_items):
    y = 1.8 + i * 0.85
    add_rounded_rect(slide, 1.5, y, 10.3, 0.7, BG_CARD, color)
    add_text_box(slide, 1.8, y + 0.05, 2.5, 0.35, title, 18, color, True)
    add_text_box(slide, 4.3, y + 0.1, 7.2, 0.5, desc, 14, LIGHT_GRAY, False)

# Stats footer
add_text_box(slide, 1.5, 7.0, 10.3, 0.3, '~370k righe  |  ~50 file JS  |  2 scenari attivi  |  0 dipendenze build-time', 14, MID_GRAY, False, PP_ALIGN.CENTER)


# ============================================================
# SALVA
# ============================================================
output_path = os.path.join(os.path.dirname(__file__), 'CVT_Panoramica_Tecnica.pptx')
prs.save(output_path)
print(f'Presentazione salvata: {output_path}')
print(f'Slide totali: {len(prs.slides)}')
