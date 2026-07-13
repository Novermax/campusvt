#!/usr/bin/env python3
"""Genera PPTX stilato come il documento originale Campus Virtual Training.pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import copy

# ===== STILI DAL TEMPLATE ORIGINALE =====
# Slide: 12192000 x 6858000 (16:9)
SLIDE_W = 12192000
SLIDE_H = 6858000

BLU_SCURO = RGBColor(0x00, 0x30, 0x63)   # #003063 titolo
BLU_MEDIO = RGBColor(0x08, 0x3C, 0x6E)   # #083C6E header slide
BLU_CHIARO = RGBColor(0x5B, 0x9B, 0xD5)  # #5B9BD5 accent1
BIANCO = RGBColor(0xFF, 0xFF, 0xFF)
NERO = RGBColor(0x00, 0x00, 0x00)
GRIGIO_CHIARO = RGBColor(0xF2, 0xF2, 0xF2)
ARANCIONE = RGBColor(0xED, 0x7D, 0x31)
GRIGIO_TESTO = RGBColor(0x44, 0x54, 0x6A)  # #44546A dk2
VERDE = RGBColor(0x70, 0xAD, 0x47)

FONT_TITOLO = 'Calibri Light'
FONT_CORPO = 'Calibri'

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def bd(slide):
    """Sfondo bianco"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BIANCO


def shape_fill(sp, color):
    """Fill solido a una shape"""
    sp.fill.solid()
    sp.fill.fore_color.rgb = color


def shape_noline(sp):
    """Rimuove il bordo"""
    ln = sp.line
    ln.fill.background()


def add_rect(slide, x, y, cx, cy, color):
    sp = slide.shapes.add_shape(1, x, y, cx, cy)
    shape_fill(sp, color)
    shape_noline(sp)
    return sp


def add_textbox(slide, x, y, cx, cy):
    return slide.shapes.add_textbox(x, y, cx, cy)


def set_para(para, text, font_name=FONT_CORPO, size=14, bold=False, color=NERO, align=PP_ALIGN.LEFT):
    para.text = text
    para.font.name = font_name
    para.font.size = Pt(size)
    para.font.bold = bold
    para.font.color.rgb = color
    para.alignment = align
    # Forza font a livello run
    for run in para.runs:
        run.font.name = font_name


def add_run(para, text, font_name=FONT_CORPO, size=14, bold=False, color=NERO):
    run = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return run


def slide_header(slide, titolo, sottotitolo=None):
    """Barra blu scuro in alto + sottotitolo"""
    add_rect(slide, 0, 0, SLIDE_W, Emu(500000), BLU_MEDIO)
    tx = add_textbox(slide, Emu(550000), Emu(60000), Emu(11000000), Emu(400000))
    set_para(tx.text_frame.paragraphs[0], titolo,
             FONT_TITOLO, 26, True, BIANCO, PP_ALIGN.LEFT)
    if sottotitolo:
        add_rect(slide, 0, Emu(500000), SLIDE_W, Emu(250000), GRIGIO_CHIARO)
        tx2 = add_textbox(slide, Emu(550000), Emu(510000), Emu(11000000), Emu(230000))
        set_para(tx2.text_frame.paragraphs[0], sottotitolo,
                 FONT_CORPO, 12, False, BLU_MEDIO, PP_ALIGN.LEFT)


def card(slide, left, top, cx, cy, titolo, corpo, color_header=BLU_MEDIO, size_corpo=12, bold_titolo=True):
    """Card con header colorato e sfondo bianco"""
    # bordo card
    sp = slide.shapes.add_shape(1, left, top, cx, cy)
    sp.fill.solid()
    sp.fill.fore_color.rgb = BIANCO
    sp.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    sp.line.width = Pt(0.5)

    # header
    hdr = add_rect(slide, left, top, cx, Emu(280000), color_header)
    tx_h = add_textbox(slide, left + Emu(80000), top + Emu(20000), cx - Emu(160000), Emu(260000))
    set_para(tx_h.text_frame.paragraphs[0], titolo,
             FONT_CORPO, 12, bold_titolo, BIANCO, PP_ALIGN.LEFT)

    # corpo
    tx_b = add_textbox(slide, left + Emu(80000), top + Emu(320000), cx - Emu(160000), cy - Emu(360000))
    tx_b.text_frame.word_wrap = True
    set_para(tx_b.text_frame.paragraphs[0], corpo,
             FONT_CORPO, size_corpo, False, NERO, PP_ALIGN.LEFT)


def make_table(slide, left, top, cx, cy, headers, rows):
    """Tabella con header blu e righe alternate"""
    nrows = len(rows) + 1
    ncols = len(headers)
    tbl = slide.shapes.add_table(nrows, ncols, left, top, cx, cy).table

    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]
        set_para(p, h, FONT_CORPO, 10, True, BIANCO, PP_ALIGN.CENTER)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLU_MEDIO
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            set_para(p, str(val), FONT_CORPO, 10, False, NERO, PP_ALIGN.CENTER)
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEF, 0xF7)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return tbl


# ====================================================================
# SLIDE 1 — TITOLO
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bd(sl)
# barra superiore
add_rect(sl, 0, 0, SLIDE_W, Emu(2500000), BLU_SCURO)
# titolo
tx = add_textbox(sl, Emu(800000), Emu(300000), Emu(10500000), Emu(1800000))
tf = tx.text_frame
tf.word_wrap = True
set_para(tf.paragraphs[0], 'Perché un cliente dovrebbe', FONT_TITOLO, 26, True, BIANCO, PP_ALIGN.CENTER)
p2 = tf.add_paragraph()
set_para(p2, 'acquistare Campus Virtual Training?', FONT_TITOLO, 26, True, BIANCO, PP_ALIGN.CENTER)
# riga arancione
add_rect(sl, Emu(4000000), Emu(1300000), Emu(4000000), Emu(120000), ARANCIONE)
# sottotitolo
tx2 = add_textbox(sl, Emu(1000000), Emu(1600000), Emu(10000000), Emu(500000))
tf2 = tx2.text_frame
tf2.word_wrap = True
set_para(tf2.paragraphs[0], 'Vantaggi economici quantificati  |  Luglio 2026', FONT_CORPO, 14, False, RGBColor(0xCC, 0xDD, 0xEE), PP_ALIGN.CENTER)

# Corpo slide
tx3 = add_textbox(sl, Emu(1500000), Emu(2800000), Emu(9000000), Emu(3500000))
tf3 = tx3.text_frame
tf3.word_wrap = True
set_para(tf3.paragraphs[0], 'La formazione tradizionale costa cara:', FONT_CORPO, 18, True, BLU_SCURO, PP_ALIGN.CENTER)
items = [
    '€800-2.000/giornata per un corso in aula SCM',
    '€300-600 di trasferta per ogni sessione',
    'Ore di fermo produzione per formare il personale',
    'Nessuna garanzia che l\'operatore abbia appreso',
]
for item in items:
    p = tf3.add_paragraph()
    set_para(p, f'✗  {item}', FONT_CORPO, 14, False, GRIGIO_TESTO, PP_ALIGN.CENTER)
p_last = tf3.add_paragraph()
set_para(p_last, '', FONT_CORPO, 10, False, NERO)
p_last = tf3.add_paragraph()
set_para(p_last, 'La risposta: CVT — tutorial 3D on-demand, nessun visore, nessuna trasferta.', FONT_CORPO, 16, True, BLU_MEDIO, PP_ALIGN.CENTER)
p_last = tf3.add_paragraph()
set_para(p_last, 'Abbonamento da €400-800/anno, attivazione in ore, disponibile 24/7.', FONT_CORPO, 16, True, BLU_MEDIO, PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 2 — I 4 VANTAGGI ECONOMICI
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bd(sl)
slide_header(sl, 'I vantaggi economici quantificati', 'Quanto vale CVT per il cliente?')

w_card = Emu(2800000)
gap = Emu(250000)
y1 = Emu(900000)
h_card = Emu(2700000)

card(sl, gap, y1, w_card, h_card,
     '1. Riduzione ticket service',
     'Un ticket SCM costa €200-800.\n'
     'Ogni tutorial evita 3-6 ticket/anno\n'
     'per deficit formativo.\n\n'
     '▶ Risparmio: €1.200/anno\n'
     '▶ Abbonamento Base: €400\n'
     '▶ ROI: 3x',
     BLU_MEDIO, 11)

card(sl, gap + w_card + gap, y1, w_card, h_card,
     '2. Riduzione downtime',
     '1h di fermo linea = €500-2.000.\n'
     'Tutorial corretto = meno errori,\n'
     'meno attese del tecnico (24-72h).\n\n'
     '▶ Risparmio: €4.000/anno\n'
     '▶ Abbonamento: €800\n'
     '▶ ROI: 5x',
     VERDE, 11)

card(sl, gap + (w_card + gap) * 2, y1, w_card, h_card,
     '3. Eliminazione corsi aula',
     'Corsi SCM: €800-2.000/giornata\n'
     '+ viaggio + hotel + ore perse.\n'
     'Sostituiti da tutorial on-demand.\n\n'
     '▶ Risparmio: €8.000-15.000/anno\n'
     '▶ Abbonamento: €400-800\n'
     '▶ ROI: 10-20x',
     ARANCIONE, 11)

card(sl, gap + (w_card + gap) * 3, y1, w_card, h_card,
     '4. Onboarding accelerato',
     'Nuovo assunto senza CVT:\n'
     '2-4 settimane per essere produttivo.\n'
     'Con CVT: 2-3 giorni.\n\n'
     '▶ Recupero: €3.000-6.000\n'
     '   per ogni nuovo assunto\n'
     '▶ Abbonamento: incluso',
     BLU_CHIARO, 11)

# Barra totale
add_rect(sl, gap, Emu(3800000), Emu(11600000), Emu(600000), BLU_SCURO)
tx_tot = add_textbox(sl, gap + Emu(100000), Emu(3850000), Emu(11400000), Emu(500000))
tf_tot = tx_tot.text_frame
tf_tot.word_wrap = True
set_para(tf_tot.paragraphs[0],
         'TOTALE risparmio annuo: €10.000 - €20.000/anno  |  Abbonamento: €400 - €800  |  ROI: 10-25x',
         FONT_CORPO, 14, True, BIANCO, PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 3 — TABELLA RIEPILOGATIVA
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bd(sl)
slide_header(sl, 'ROI complessivo per il cliente')

make_table(sl, Emu(400000), Emu(900000), Emu(11400000), Emu(2000000),
    ['Voce di risparmio', 'Risparmio annuo', 'Abbonamento', 'ROI'],
    [
        ['Ticket service evitati (3-5 ticket/anno)', '€1.200', '€400 (Base)', '3x'],
        ['Downtime ridotto (4h/anno evitate)', '€4.000', '€800 (Professional)', '5x'],
        ['Corsi aula eliminati', '€8.000-15.000', '€400-800', '10-20x'],
        ['Onboarding accelerato (per assunto)', '€3.000-6.000', 'Incluso', '—'],
    ])

tx = add_textbox(sl, Emu(400000), Emu(3200000), Emu(11400000), Emu(1800000))
tf = tx.text_frame
tf.word_wrap = True
msg = (
    'Il primo tutorial seguito (es. manutenzione pompa Becker)\n'
    'evita già almeno 1 chiamata al service — €200-800.\n\n'
    'L\'investimento è ripagato nel primo mese. Ogni tutorial successivo è guadagno netto.'
)
set_para(tf.paragraphs[0], msg, FONT_CORPO, 16, True, BLU_MEDIO, PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 4 — PERCHÉ SI RINNOVA
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bd(sl)
slide_header(sl, 'Perché il rinnovo è automatico', 'Vantaggi che crescono nel tempo')

items = [
    ('Catalogo crescente', 'Ogni anno arrivano tutorial per\nnuove macchine e procedure.\nPiù valore, stesso prezzo.'),
    ('Sostituzione manuali', 'I manuali cartacei vengono\neliminati. Il personale si abitua\nal formato 3D interattivo.'),
    ('Onboarding infrastrutturale', 'Nuovi assunti formati in giorni.\nLa piattaforma diventa\ninfrastruttura aziendale.'),
    ('Tutorial proprietari', 'Il cliente crea i propri contenuti\ncon l\'editor. Il valore cresce\ncon l\'uso.'),
    ('Lock-in positivo', 'Più tutorial si accumulano,\npiù la piattaforma è\nindispensabile.'),
]

w = Emu(2100000)
g = Emu(200000)
y = Emu(1000000)
h = Emu(3000000)

for i, (tit, corpo) in enumerate(items):
    card(sl, g + (w + g) * i, y, w, h, tit, corpo,
         BLU_MEDIO if i % 2 == 0 else BLU_CHIARO, 11)

tx = add_textbox(sl, Emu(400000), Emu(4300000), Emu(11400000), Emu(600000))
tf = tx.text_frame
tf.word_wrap = True
set_para(tf.paragraphs[0],
         'Anche se il prezzo raddoppiasse, CVT costerebbe ancora 10x meno dei competitor enterprise.',
         FONT_CORPO, 14, True, GRIGIO_TESTO, PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 5 — CONFRONTO COMPETITOR
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bd(sl)
slide_header(sl, 'Confronto con le alternative enterprise', 'Perché non comprare Strivr, Talespin o EON?')

make_table(sl, Emu(300000), Emu(900000), Emu(11600000), Emu(1800000),
    ['Prodotto', 'Prezzo annuo', 'VR?', 'Tempi', 'Target'],
    [
        ['Campus Virtual Training', '€400 - €5.000', 'No', 'Ore', 'PMI (95% mercato)'],
        ['Strivr', '$50.000 - $250.000', 'Sì', 'Mesi', 'Fortune 500'],
        ['Talespin', '$30.000 - $150.000', 'Sì', 'Mesi', 'Grandi aziende'],
        ['EON Reality', '€10.000 - €50.000', 'Opz.', 'Settimane', 'Education + Enterprise'],
        ['PTC Vuforia', '$5.000 - $100.000+', 'Opz.', 'Settimane', 'Grandi manifatturiere'],
    ])

tx = add_textbox(sl, Emu(300000), Emu(3000000), Emu(11600000), Emu(2000000))
tf = tx.text_frame
tf.word_wrap = True
set_para(tf.paragraphs[0], 'Il 95% del tessuto industriale italiano sono PMI con 5-50 dipendenti.', FONT_CORPO, 16, True, BLU_SCURO, PP_ALIGN.CENTER)
p = tf.add_paragraph()
set_para(p, '', FONT_CORPO, 8)
items = [
    'Strivr/Talespin: budget annuale superiore all\'intero budget formazione di una PMI',
    'Richiedono headset VR (€500-1.000/cad.) + consulenza implementativa',
    'Tempi di attivazione in mesi, quando serve subito',
]
for item in items:
    p = tf.add_paragraph()
    set_para(p, f'✗  {item}', FONT_CORPO, 13, False, GRIGIO_TESTO)
p = tf.add_paragraph()
set_para(p, '', FONT_CORPO, 8)
p = tf.add_paragraph()
set_para(p, 'CVT è l\'unica opzione accessibile per la formazione industriale 3D.', FONT_CORPO, 16, True, VERDE, PP_ALIGN.CENTER)

# ====================================================================
# SLIDE 6 — FRASE SINTESI FINALE
# ====================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bd(sl)
# barra piena
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, BLU_SCURO)

tx = add_textbox(sl, Emu(800000), Emu(1000000), Emu(10600000), Emu(2400000))
tf = tx.text_frame
tf.word_wrap = True
set_para(tf.paragraphs[0], '"Con €400/anno', FONT_TITOLO, 32, True, BIANCO, PP_ALIGN.CENTER)
p = tf.add_paragraph()
set_para(p, '', FONT_CORPO, 12)
p = tf.add_paragraph()
set_para(p, 'eviti almeno 1 ticket service (€200-800)', FONT_TITOLO, 28, True, ARANCIONE, PP_ALIGN.CENTER)
p = tf.add_paragraph()
set_para(p, 'e riduci il downtime produzione."', FONT_TITOLO, 28, True, BIANCO, PP_ALIGN.CENTER)

tx2 = add_textbox(sl, Emu(1200000), Emu(3600000), Emu(9800000), Emu(2000000))
tf2 = tx2.text_frame
tf2.word_wrap = True
set_para(tf2.paragraphs[0], '10 operatori × 5 tutorial = ROI > 500% annuo', FONT_CORPO, 18, True, RGBColor(0xAA, 0xCC, 0xEE), PP_ALIGN.CENTER)
p = tf2.add_paragraph()
set_para(p, 'Il primo mese ripaga l\'intero abbonamento.', FONT_CORPO, 18, True, RGBColor(0xCC, 0xDD, 0xFF), PP_ALIGN.CENTER)

tx3 = add_textbox(sl, Emu(1000000), Emu(5800000), Emu(10200000), Emu(500000))
tf3 = tx3.text_frame
set_para(tf3.paragraphs[0], 'Campus Virtual Training  ■  Analisi strategica  ■  Luglio 2026', FONT_CORPO, 11, False, RGBColor(0x88, 0xAA, 0xCC), PP_ALIGN.CENTER)

# ===== SALVA =====
out = r'C:\Users\mloffredo\campus_virtual_training\CVT_Value_Proposition.pptx'
prs.save(out)
print(f'Salvato: {out}')
print('Fatto!')
